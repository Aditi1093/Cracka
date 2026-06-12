"""
utils/office_creator.py
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
CRACKA AI — Office Document Creator
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Features:
  ✅ PowerPoint — AI se content generate karke slides banao
  ✅ Word       — Formatted documents, reports, resumes
  ✅ Excel      — Data sheets, budgets, trackers

Voice Commands:
  "make presentation on [topic]"
  "create ppt on [topic]"
  "make slides on [topic]"

  "write word document on [topic]"
  "create report on [topic]"
  "make resume"

  "create excel sheet for [purpose]"
  "make budget sheet"
  "create tracker for [purpose]"

All files save to Desktop automatically and open in respective apps.
"""

import os
import re
import subprocess
from datetime import datetime
from brain.chat_engine import ask_ai
from core.logger import log_info, log_error
from core.voice_engine import speak

DESKTOP = os.path.join(os.path.expanduser("~"), "Desktop")
os.makedirs(DESKTOP, exist_ok=True)


def _safe_filename(name: str) -> str:
    """Convert topic to safe filename."""
    name = re.sub(r'[<>:"/\\|?*]', '', name)
    name = name.replace(' ', '_')[:40]
    ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{name}_{ts}"


def _open_file(filepath: str):
    """Open file with default application."""
    try:
        os.startfile(filepath)
    except Exception:
        try:
            subprocess.Popen(["start", "", filepath], shell=True)
        except Exception as e:
            log_error(f"Could not open file: {e}")


# ══════════════════════════════════════════════════════════════
# POWERPOINT CREATOR
# ══════════════════════════════════════════════════════════════

def create_presentation(command: str) -> str:
    """
    Create a PowerPoint presentation on any topic.
    Voice: "make presentation on AI" / "create ppt on climate change"
    """
    # Extract topic
    topic = command.lower()
    for phrase in [
        "make presentation on", "create presentation on",
        "make ppt on", "create ppt on", "make slides on",
        "create slides on", "presentation on", "ppt on",
        "slides on", "make a presentation on", "create a presentation on",
        "make a ppt on", "presentation about", "slides about",
    ]:
        topic = topic.replace(phrase, "").strip()

    topic = topic.strip(" .")
    if not topic:
        return "Please tell me the topic Boss. Like 'make presentation on AI'."

    speak(f"Creating presentation on {topic} Boss. Give me a moment.")
    log_info(f"Creating PPT: {topic}")

    # Generate content with AI
    prompt = f"""Create a professional PowerPoint presentation on: {topic}

Generate exactly 7 slides with this JSON structure (respond ONLY with JSON, no other text):

{{
  "title": "Main presentation title",
  "subtitle": "Subtitle or tagline",
  "slides": [
    {{
      "slide_number": 1,
      "title": "Slide Title",
      "content": ["Point 1", "Point 2", "Point 3", "Point 4"],
      "notes": "Speaker notes for this slide"
    }}
  ]
}}

Slides should be:
1. Title slide
2. Introduction / Overview
3. Main concept 1
4. Main concept 2
5. Key statistics or facts
6. Applications / Use cases
7. Conclusion / Summary

Make content informative and professional."""

    try:
        raw = ask_ai(prompt)
        # Extract JSON
        json_match = re.search(r'\{.*\}', raw, re.DOTALL)
        if json_match:
            import json
            data = json.loads(json_match.group())
        else:
            raise ValueError("No JSON found in response")
    except Exception as e:
        log_error(f"AI content generation error: {e}")
        # Fallback content
        data = _get_fallback_ppt_content(topic)

    return _build_pptx(data, topic)


def _get_fallback_ppt_content(topic: str) -> dict:
    """Fallback PPT content if AI fails."""
    return {
        "title": topic.title(),
        "subtitle": f"A Comprehensive Overview",
        "slides": [
            {
                "slide_number": 1,
                "title": topic.title(),
                "content": [f"Comprehensive overview of {topic}", "Created by Cracka AI"],
                "notes": "Title slide"
            },
            {
                "slide_number": 2,
                "title": "Introduction",
                "content": [
                    f"What is {topic}?",
                    "Why it matters",
                    "Current trends",
                    "Future outlook"
                ],
                "notes": "Introduction slide"
            },
            {
                "slide_number": 3,
                "title": "Key Concepts",
                "content": [
                    "Fundamental principles",
                    "Core components",
                    "Important definitions",
                    "Basic framework"
                ],
                "notes": "Key concepts"
            },
            {
                "slide_number": 4,
                "title": "Applications",
                "content": [
                    "Real-world use cases",
                    "Industry applications",
                    "Examples and case studies",
                    "Impact assessment"
                ],
                "notes": "Applications slide"
            },
            {
                "slide_number": 5,
                "title": "Benefits & Challenges",
                "content": [
                    "Key advantages",
                    "Main benefits",
                    "Potential challenges",
                    "Solutions and mitigations"
                ],
                "notes": "Benefits and challenges"
            },
            {
                "slide_number": 6,
                "title": "Future Trends",
                "content": [
                    "Emerging developments",
                    "Future possibilities",
                    "Research directions",
                    "Innovation opportunities"
                ],
                "notes": "Future trends"
            },
            {
                "slide_number": 7,
                "title": "Conclusion",
                "content": [
                    "Key takeaways",
                    "Summary of main points",
                    "Next steps",
                    "Thank you"
                ],
                "notes": "Conclusion slide"
            }
        ]
    }


def _build_pptx(data: dict, topic: str) -> str:
    """Build the actual PowerPoint file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # ── Color scheme ──────────────────────────────────────────────────────
        DARK_BG    = RGBColor(0x0A, 0x0A, 0x1A)   # Dark navy
        CYAN       = RGBColor(0x00, 0xE5, 0xFF)   # Cyan accent
        WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
        ACCENT     = RGBColor(0x00, 0xC8, 0xD4)

        slides_data = data.get("slides", [])

        for i, slide_data in enumerate(slides_data):
            slide_title   = slide_data.get("title", f"Slide {i+1}")
            slide_content = slide_data.get("content", [])
            slide_notes   = slide_data.get("notes", "")

            # Choose layout based on slide number
            if i == 0:
                # Title slide — blank layout
                slide_layout = prs.slide_layouts[6]
                slide        = prs.slides.add_slide(slide_layout)

                # Dark background
                background = slide.background
                fill       = background.fill
                fill.solid()
                fill.fore_color.rgb = DARK_BG

                # Cyan accent bar (top)
                bar_top      = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(13.33), Inches(0.12)
                )
                bar_top.fill.solid()
                bar_top.fill.fore_color.rgb = CYAN
                bar_top.line.fill.background()

                # Main title
                title_box = slide.shapes.add_textbox(
                    Inches(1), Inches(2.0), Inches(11.33), Inches(2)
                )
                tf   = title_box.text_frame
                tf.word_wrap = True
                p    = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run  = p.add_run()
                run.text = data.get("title", slide_title).upper()
                run.font.size  = Pt(44)
                run.font.bold  = True
                run.font.color.rgb = CYAN
                run.font.name  = "Consolas"

                # Subtitle
                sub_box = slide.shapes.add_textbox(
                    Inches(1), Inches(4.2), Inches(11.33), Inches(1)
                )
                tf2  = sub_box.text_frame
                p2   = tf2.paragraphs[0]
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = data.get("subtitle", "")
                run2.font.size  = Pt(20)
                run2.font.color.rgb = LIGHT_GRAY
                run2.font.name  = "Consolas"

                # Bottom bar
                bar_bot      = slide.shapes.add_shape(
                    1, Inches(0), Inches(7.3), Inches(13.33), Inches(0.12)
                )
                bar_bot.fill.solid()
                bar_bot.fill.fore_color.rgb = CYAN
                bar_bot.line.fill.background()

                # Cracka watermark
                wm_box = slide.shapes.add_textbox(
                    Inches(9), Inches(6.9), Inches(4), Inches(0.4)
                )
                wm_tf  = wm_box.text_frame
                wm_p   = wm_tf.paragraphs[0]
                wm_p.alignment = PP_ALIGN.RIGHT
                wm_run = wm_p.add_run()
                wm_run.text = "Created by CRACKA AI"
                wm_run.font.size  = Pt(9)
                wm_run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
                wm_run.font.name  = "Consolas"

            else:
                # Content slide
                slide_layout = prs.slide_layouts[6]
                slide        = prs.slides.add_slide(slide_layout)

                # Background
                background = slide.background
                fill       = background.fill
                fill.solid()
                fill.fore_color.rgb = DARK_BG

                # Top accent bar
                top_bar = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(13.33), Inches(0.08)
                )
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = CYAN
                top_bar.line.fill.background()

                # Slide number badge
                num_box = slide.shapes.add_shape(
                    1, Inches(0), Inches(0.08), Inches(0.5), Inches(0.55)
                )
                num_box.fill.solid()
                num_box.fill.fore_color.rgb = CYAN
                num_box.line.fill.background()
                num_tf  = num_box.text_frame
                num_p   = num_tf.paragraphs[0]
                num_p.alignment = PP_ALIGN.CENTER
                num_run = num_p.add_run()
                num_run.text = str(i + 1)
                num_run.font.size  = Pt(12)
                num_run.font.bold  = True
                num_run.font.color.rgb = DARK_BG
                num_run.font.name  = "Consolas"

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.6), Inches(0.1), Inches(12.5), Inches(0.7)
                )
                title_tf  = title_box.text_frame
                title_p   = title_tf.paragraphs[0]
                title_run = title_p.add_run()
                title_run.text = slide_title.upper()
                title_run.font.size  = Pt(26)
                title_run.font.bold  = True
                title_run.font.color.rgb = CYAN
                title_run.font.name  = "Consolas"

                # Divider line
                div = slide.shapes.add_shape(
                    1, Inches(0.6), Inches(0.85), Inches(12), Inches(0.03)
                )
                div.fill.solid()
                div.fill.fore_color.rgb = ACCENT
                div.line.fill.background()

                # Content bullets
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.1), Inches(11.5), Inches(5.8)
                )
                content_tf        = content_box.text_frame
                content_tf.word_wrap = True

                for j, point in enumerate(slide_content):
                    if j == 0:
                        cp = content_tf.paragraphs[0]
                    else:
                        cp = content_tf.add_paragraph()

                    cp.space_before = Pt(8)
                    cp.space_after  = Pt(4)

                    # Bullet dot
                    dot_run = cp.add_run()
                    dot_run.text = "◉  "
                    dot_run.font.size  = Pt(10)
                    dot_run.font.color.rgb = CYAN
                    dot_run.font.name  = "Consolas"

                    # Content text
                    text_run = cp.add_run()
                    text_run.text = str(point)
                    text_run.font.size  = Pt(18)
                    text_run.font.color.rgb = WHITE
                    text_run.font.name  = "Calibri"

                # Bottom bar
                bot_bar = slide.shapes.add_shape(
                    1, Inches(0), Inches(7.3), Inches(13.33), Inches(0.08)
                )
                bot_bar.fill.solid()
                bot_bar.fill.fore_color.rgb = CYAN
                bot_bar.line.fill.background()

                # Topic label bottom right
                topic_box = slide.shapes.add_textbox(
                    Inches(9), Inches(7.1), Inches(4.2), Inches(0.3)
                )
                topic_tf  = topic_box.text_frame
                topic_p   = topic_tf.paragraphs[0]
                topic_p.alignment = PP_ALIGN.RIGHT
                topic_run = topic_p.add_run()
                topic_run.text = f"CRACKA AI  |  {topic.upper()[:30]}"
                topic_run.font.size  = Pt(8)
                topic_run.font.color.rgb = RGBColor(0x33, 0x33, 0x55)
                topic_run.font.name  = "Consolas"

            # Add speaker notes
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_tf    = notes_slide.notes_text_frame
                notes_tf.text = slide_notes

        # Save file
        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}.pptx")
        prs.save(filepath)

        log_info(f"PPT saved: {filepath}")
        _open_file(filepath)

        speak(f"Presentation on {topic} is ready Boss! Opening now.")
        return (f"Presentation created Boss! "
                f"{len(slides_data)} slides on {topic}. "
                f"Saved to Desktop and opening now.")

    except ImportError:
        return "python-pptx not installed Boss. Run: pip install python-pptx"
    except Exception as e:
        log_error(f"PPT creation error: {e}")
        return f"Could not create presentation Boss. Error: {e}"


# ══════════════════════════════════════════════════════════════
# WORD DOCUMENT CREATOR
# ══════════════════════════════════════════════════════════════

def create_word_document(command: str) -> str:
    """
    Create a Word document on any topic.
    Voice: "write word document on cybersecurity"
           "create report on machine learning"
           "make resume" / "write essay on AI"
    """
    cmd_lower = command.lower()

    # Detect document type
    is_resume   = "resume" in cmd_lower or "cv" in cmd_lower
    is_report   = "report" in cmd_lower
    is_essay    = "essay" in cmd_lower
    is_letter   = "letter" in cmd_lower
    is_proposal = "proposal" in cmd_lower

    # Extract topic
    topic = cmd_lower
    for phrase in [
        "write word document on", "create word document on",
        "make word document on", "write document on",
        "create report on", "write report on",
        "write essay on", "create essay on",
        "write a report on", "create a report on",
        "make a report on", "word document on",
        "write letter to", "create proposal for",
        "make resume", "create resume", "write resume",
        "make cv", "create cv",
    ]:
        topic = topic.replace(phrase, "").strip()

    topic = topic.strip(" .")

    if is_resume:
        return _create_resume()

    if not topic:
        return "Please tell me the topic Boss. Like 'write report on cybersecurity'."

    speak(f"Writing document on {topic} Boss.")
    log_info(f"Creating Word doc: {topic}")

    # Generate content with AI
    if is_report:
        prompt = f"""Write a professional report on: {topic}

Structure:
1. Executive Summary (2-3 sentences)
2. Introduction (1 paragraph)
3. Main Content (3-4 sections with headings)
4. Key Findings (bullet points)
5. Recommendations (3-5 points)
6. Conclusion (1 paragraph)

Make it professional and informative. About 500-700 words total."""
    elif is_essay:
        prompt = f"""Write a well-structured essay on: {topic}

Structure:
- Introduction (hook + thesis)
- Body paragraph 1 (main argument)
- Body paragraph 2 (supporting evidence)
- Body paragraph 3 (counter-argument and rebuttal)
- Conclusion (summary + broader implications)

About 400-600 words."""
    else:
        prompt = f"""Write a comprehensive document on: {topic}

Include:
- Introduction
- Key concepts and definitions
- Main topics (3-4 sections)
- Important points in each section
- Conclusion

Professional tone. About 400-600 words."""

    content = ask_ai(prompt)

    return _build_docx(content, topic, is_report, is_essay)


def _create_resume() -> str:
    """Create a professional resume template."""
    speak("Creating your resume template Boss.")

    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()

        # Page margins
        for section in doc.sections:
            section.top_margin    = Inches(0.7)
            section.bottom_margin = Inches(0.7)
            section.left_margin   = Inches(0.8)
            section.right_margin  = Inches(0.8)

        # Helper functions
        def add_heading(text, level=1, color_hex="000000"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(text)
            run.bold = True
            run.font.size = Pt(14 if level == 1 else 12)
            run.font.color.rgb = RGBColor(
                int(color_hex[0:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16)
            )
            return p

        def add_line(color_hex="0077AA"):
            p   = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pBdr= OxmlElement('w:pBdr')
            bottom = OxmlElement('w:bottom')
            bottom.set(qn('w:val'), 'single')
            bottom.set(qn('w:sz'), '8')
            bottom.set(qn('w:space'), '1')
            bottom.set(qn('w:color'), color_hex)
            pBdr.append(bottom)
            pPr.append(pBdr)
            p.paragraph_format.space_after  = Pt(2)
            p.paragraph_format.space_before = Pt(2)

        # ── NAME ──────────────────────────────────────────────────────────────
        name_p = doc.add_paragraph()
        name_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        name_run = name_p.add_run("YOUR FULL NAME")
        name_run.bold = True
        name_run.font.size = Pt(24)
        name_run.font.color.rgb = RGBColor(0x00, 0x77, 0xAA)

        # Contact info
        contact_p = doc.add_paragraph()
        contact_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        contact_run = contact_p.add_run(
            "📧 your.email@gmail.com  |  📱 +91 XXXXX XXXXX  |  "
            "🔗 linkedin.com/in/yourname  |  💻 github.com/yourname"
        )
        contact_run.font.size = Pt(10)
        contact_run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        add_line()

        # ── OBJECTIVE ─────────────────────────────────────────────────────────
        add_heading("CAREER OBJECTIVE", color_hex="0077AA")
        add_line()
        obj_p = doc.add_paragraph()
        obj_p.add_run(
            "Motivated Computer Science Engineering student specializing in "
            "Cybersecurity, seeking opportunities to apply AI and security skills "
            "to solve real-world problems and contribute to innovative projects."
        ).font.size = Pt(10)

        # ── EDUCATION ─────────────────────────────────────────────────────────
        doc.add_paragraph()
        add_heading("EDUCATION", color_hex="0077AA")
        add_line()

        edu_table = doc.add_table(rows=2, cols=3)
        edu_table.style = "Table Grid"
        # Header
        hdr_cells = edu_table.rows[0].cells
        for cell, text in zip(hdr_cells, ["Degree", "Institution", "Year/CGPA"]):
            cell.paragraphs[0].add_run(text).bold = True
            cell.paragraphs[0].runs[0].font.size = Pt(10)
        # Data
        row_cells = edu_table.rows[1].cells
        for cell, text in zip(row_cells, [
            "B.Tech CSE (Cybersecurity)",
            "Your University Name",
            "2023-2027 | CGPA: X.X"
        ]):
            cell.paragraphs[0].add_run(text).font.size = Pt(10)

        # ── SKILLS ────────────────────────────────────────────────────────────
        doc.add_paragraph()
        add_heading("TECHNICAL SKILLS", color_hex="0077AA")
        add_line()

        skills = {
            "Programming Languages": "Python, C/C++, JavaScript",
            "AI / ML":               "TensorFlow, PyTorch, Scikit-learn, LangChain",
            "Cybersecurity":         "Network Security, Ethical Hacking, OWASP, Wireshark, Burp Suite",
            "Tools & Technologies":  "Git, Linux, Docker, VS Code, Ollama",
            "Frameworks":            "FastAPI, Flask, PyQt5, OpenCV",
        }
        for skill, detail in skills.items():
            p = doc.add_paragraph(style="List Bullet")
            run1 = p.add_run(f"{skill}: ")
            run1.bold = True
            run1.font.size = Pt(10)
            run2 = p.add_run(detail)
            run2.font.size = Pt(10)

        # ── PROJECTS ──────────────────────────────────────────────────────────
        doc.add_paragraph()
        add_heading("PROJECTS", color_hex="0077AA")
        add_line()

        projects = [
            {
                "name": "Cracka AI — Personal AI Assistant",
                "tech": "Python, PyQt5, Ollama, Vosk, OpenCV",
                "desc": [
                    "Built a complete AI voice assistant with 30+ features",
                    "Integrated cyber security tools: network monitor, port scanner, phishing detector",
                    "Implemented wake word detection, multilingual support, face recognition",
                ]
            },
            {
                "name": "Project 2 Name",
                "tech": "Technologies used",
                "desc": [
                    "Key achievement 1",
                    "Key achievement 2",
                ]
            },
        ]

        for proj in projects:
            p = doc.add_paragraph()
            name_run = p.add_run(f"• {proj['name']}")
            name_run.bold = True
            name_run.font.size = Pt(11)
            name_run.font.color.rgb = RGBColor(0x00, 0x77, 0xAA)

            tech_p = doc.add_paragraph()
            tech_p.paragraph_format.left_indent = Inches(0.2)
            tech_run = tech_p.add_run(f"Tech Stack: {proj['tech']}")
            tech_run.italic = True
            tech_run.font.size = Pt(9)
            tech_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

            for desc in proj["desc"]:
                dp = doc.add_paragraph(style="List Bullet")
                dp.paragraph_format.left_indent = Inches(0.3)
                dp.add_run(desc).font.size = Pt(10)

        # ── CERTIFICATIONS ────────────────────────────────────────────────────
        doc.add_paragraph()
        add_heading("CERTIFICATIONS & ACHIEVEMENTS", color_hex="0077AA")
        add_line()

        certs = [
            "TryHackMe — Jr Penetration Tester Path (In Progress)",
            "Python Programming — [Platform Name]",
            "Add your certifications here",
        ]
        for cert in certs:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(cert).font.size = Pt(10)

        # ── LANGUAGES ─────────────────────────────────────────────────────────
        doc.add_paragraph()
        add_heading("LANGUAGES", color_hex="0077AA")
        add_line()
        lang_p = doc.add_paragraph()
        lang_p.add_run("English (Professional)  •  Hindi (Native)  •  Marathi (Native)").font.size = Pt(10)

        # Save
        filename = f"Resume_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        filepath = os.path.join(DESKTOP, f"{filename}.docx")
        doc.save(filepath)

        log_info(f"Resume saved: {filepath}")
        _open_file(filepath)

        speak("Resume template created Boss! Opening in Word now. Fill in your details.")
        return f"Resume created Boss! Opening on Desktop. Fill in your details."

    except ImportError:
        return "python-docx not installed Boss. Run: pip install python-docx"
    except Exception as e:
        log_error(f"Resume creation error: {e}")
        return f"Could not create resume Boss: {e}"


def _build_docx(content: str, topic: str,
                is_report: bool = False, is_essay: bool = False) -> str:
    """Build Word document from AI content."""
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()

        # Page setup
        for section in doc.sections:
            section.top_margin    = Inches(1.0)
            section.bottom_margin = Inches(1.0)
            section.left_margin   = Inches(1.2)
            section.right_margin  = Inches(1.2)

        # ── Title ─────────────────────────────────────────────────────────────
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_p.add_run(topic.title())
        title_run.bold = True
        title_run.font.size = Pt(22)
        title_run.font.color.rgb = RGBColor(0x00, 0x77, 0xAA)
        title_run.font.name = "Calibri"

        # Document type label
        doc_type = "Report" if is_report else ("Essay" if is_essay else "Document")
        type_p = doc.add_paragraph()
        type_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        type_run = type_p.add_run(
            f"{doc_type}  |  Generated by Cracka AI  |  "
            f"{datetime.now().strftime('%B %d, %Y')}"
        )
        type_run.font.size  = Pt(10)
        type_run.italic     = True
        type_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        doc.add_paragraph()  # Spacer

        # ── Content ───────────────────────────────────────────────────────────
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Detect headings (lines starting with # or ALL CAPS short lines)
            if line.startswith('# '):
                h = doc.add_heading(line[2:], level=1)
                h.runs[0].font.color.rgb = RGBColor(0x00, 0x77, 0xAA)

            elif line.startswith('## '):
                h = doc.add_heading(line[3:], level=2)
                h.runs[0].font.color.rgb = RGBColor(0x00, 0x99, 0xCC)

            elif line.startswith('### '):
                h = doc.add_heading(line[4:], level=3)

            # Bullet points
            elif line.startswith(('- ', '• ', '* ')):
                p = doc.add_paragraph(style='List Bullet')
                p.add_run(line[2:]).font.size = Pt(11)

            # Numbered lists
            elif re.match(r'^\d+\.\s', line):
                p = doc.add_paragraph(style='List Number')
                p.add_run(re.sub(r'^\d+\.\s', '', line)).font.size = Pt(11)

            # Bold lines (markdown **text**)
            elif line.startswith('**') and line.endswith('**'):
                p = doc.add_paragraph()
                run = p.add_run(line[2:-2])
                run.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x00, 0x77, 0xAA)

            # Regular paragraph
            else:
                p = doc.add_paragraph()
                run = p.add_run(line)
                run.font.size = Pt(11)
                run.font.name = "Calibri"
                p.paragraph_format.space_after  = Pt(6)
                p.paragraph_format.line_spacing  = Pt(14)

        # Save
        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}.docx")
        doc.save(filepath)

        log_info(f"Word doc saved: {filepath}")
        _open_file(filepath)

        speak(f"Document on {topic} is ready Boss! Opening in Word now.")
        return f"Word document created Boss! Saved to Desktop and opening now."

    except ImportError:
        return "python-docx not installed Boss. Run: pip install python-docx"
    except Exception as e:
        log_error(f"Word doc error: {e}")
        return f"Could not create document Boss: {e}"


# ══════════════════════════════════════════════════════════════
# EXCEL CREATOR
# ══════════════════════════════════════════════════════════════

def create_excel_sheet(command: str) -> str:
    """
    Create Excel spreadsheet.
    Voice: "create excel sheet for budget"
           "make tracker for expenses"
           "create data sheet for students"
    """
    cmd_lower = command.lower()

    # Detect type
    is_budget   = any(w in cmd_lower for w in ["budget", "expense", "finance", "money", "cost"])
    is_tracker  = any(w in cmd_lower for w in ["tracker", "track", "habit", "goal", "progress"])
    is_schedule = any(w in cmd_lower for w in ["schedule", "timetable", "calendar", "planner"])
    is_student  = any(w in cmd_lower for w in ["student", "marks", "grades", "attendance"])
    is_inventory= any(w in cmd_lower for w in ["inventory", "stock", "product", "item"])

    # Extract topic
    topic = cmd_lower
    for phrase in [
        "create excel sheet for", "make excel sheet for",
        "create excel for", "make excel for",
        "create spreadsheet for", "make spreadsheet for",
        "create tracker for", "make tracker for",
        "create sheet for", "make sheet for",
        "excel sheet for", "excel for",
    ]:
        topic = topic.replace(phrase, "").strip()
    topic = topic.strip(" .")

    if not topic:
        topic = "data"

    speak(f"Creating Excel sheet for {topic} Boss.")

    if is_budget:
        return _create_budget_excel(topic)
    elif is_tracker:
        return _create_tracker_excel(topic)
    elif is_schedule:
        return _create_schedule_excel(topic)
    elif is_student:
        return _create_student_excel(topic)
    else:
        return _create_generic_excel(topic)


def _apply_excel_styles(ws, header_cells, header_color="0077AA"):
    """Apply professional styling to Excel sheet."""
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side

    header_fill = PatternFill("solid", fgColor=header_color.replace("#", ""))
    border_side = Side(style='thin', color='CCCCCC')
    border      = Border(
        left=border_side, right=border_side,
        top=border_side,  bottom=border_side
    )

    for cell in header_cells:
        cell.fill      = header_fill
        cell.font      = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border    = border

    return header_fill, border


def _create_budget_excel(topic: str) -> str:
    """Create a budget/expense tracker Excel sheet."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()

        # ── Sheet 1: Monthly Budget ───────────────────────────────────────────
        ws1 = wb.active
        ws1.title = "Monthly Budget"

        months = ["January", "February", "March", "April", "May", "June",
                  "July", "August", "September", "October", "November", "December"]

        categories = [
            "Rent/Housing", "Food & Groceries", "Transportation",
            "Utilities (Electric/Water)", "Internet & Phone",
            "Entertainment", "Clothing", "Healthcare",
            "Education", "Savings", "Miscellaneous"
        ]

        # Title
        ws1.merge_cells("A1:N1")
        ws1["A1"] = f"MONTHLY BUDGET TRACKER — {topic.upper()}"
        ws1["A1"].font = Font(bold=True, size=14, color="FFFFFF", name="Calibri")
        ws1["A1"].fill = PatternFill("solid", fgColor="0077AA")
        ws1["A1"].alignment = Alignment(horizontal="center")
        ws1.row_dimensions[1].height = 25

        # Headers
        ws1["A2"] = "Category"
        for i, month in enumerate(months[:6], 2):
            ws1.cell(row=2, column=i).value = month
        ws1.cell(row=2, column=8).value = "H1 Total"

        header_cells = [ws1.cell(row=2, column=c) for c in range(1, 9)]
        _apply_excel_styles(ws1, header_cells)

        # Category rows
        for r, cat in enumerate(categories, 3):
            ws1.cell(row=r, column=1).value = cat
            ws1.cell(row=r, column=1).font = Font(name="Calibri", size=10)

            for c in range(2, 8):
                cell = ws1.cell(row=r, column=c)
                cell.value = 0
                cell.number_format = '₹#,##0.00'
                border_side = Side(style='thin', color='CCCCCC')
                cell.border = Border(
                    left=border_side, right=border_side,
                    top=border_side, bottom=border_side
                )

            # H1 Total formula
            total_cell = ws1.cell(row=r, column=8)
            col_letters = [get_column_letter(c) for c in range(2, 8)]
            formula = "+".join([f"{col}{r}" for col in col_letters])
            total_cell.value = f"={formula}"
            total_cell.number_format = '₹#,##0.00'
            total_cell.font = Font(bold=True, name="Calibri", size=10)

        # Total row
        total_row = len(categories) + 3
        ws1.cell(row=total_row, column=1).value = "TOTAL"
        ws1.cell(row=total_row, column=1).font = Font(bold=True, color="FFFFFF")
        ws1.cell(row=total_row, column=1).fill = PatternFill("solid", fgColor="0077AA")

        for c in range(2, 9):
            col_letter = get_column_letter(c)
            cell = ws1.cell(row=total_row, column=c)
            cell.value = f"=SUM({col_letter}3:{col_letter}{total_row-1})"
            cell.number_format = '₹#,##0.00'
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0077AA")

        # Column widths
        ws1.column_dimensions["A"].width = 22
        for c in range(2, 9):
            ws1.column_dimensions[get_column_letter(c)].width = 14

        # ── Sheet 2: Expense Log ──────────────────────────────────────────────
        ws2 = wb.create_sheet("Expense Log")

        ws2.merge_cells("A1:F1")
        ws2["A1"] = "DAILY EXPENSE LOG"
        ws2["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws2["A1"].fill = PatternFill("solid", fgColor="0077AA")
        ws2["A1"].alignment = Alignment(horizontal="center")

        headers2 = ["Date", "Category", "Description", "Amount (₹)", "Payment Mode", "Notes"]
        for c, h in enumerate(headers2, 1):
            cell = ws2.cell(row=2, column=c)
            cell.value = h
        _apply_excel_styles(ws2, [ws2.cell(row=2, column=c) for c in range(1, 7)])

        # Sample rows
        sample_data = [
            ["01/06/2025", "Food & Groceries", "Big Bazaar shopping", 1500, "UPI", "Monthly groceries"],
            ["02/06/2025", "Transportation",   "Auto rickshaw",        80,  "Cash", "Office commute"],
            ["03/06/2025", "Entertainment",    "Netflix subscription", 649, "Card", "Monthly subscription"],
        ]
        for r, row in enumerate(sample_data, 3):
            for c, val in enumerate(row, 1):
                ws2.cell(row=r, column=c).value = val
                if c == 4:
                    ws2.cell(row=r, column=c).number_format = '₹#,##0.00'

        col_widths2 = [12, 20, 25, 15, 15, 20]
        for c, w in enumerate(col_widths2, 1):
            ws2.column_dimensions[get_column_letter(c)].width = w

        # ── Sheet 3: Summary Chart Data ────────────────────────────────────────
        ws3 = wb.create_sheet("Summary")
        ws3.merge_cells("A1:C1")
        ws3["A1"] = "BUDGET SUMMARY"
        ws3["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws3["A1"].fill = PatternFill("solid", fgColor="0077AA")
        ws3["A1"].alignment = Alignment(horizontal="center")

        ws3["A2"] = "Category"
        ws3["B2"] = "Budgeted (₹)"
        ws3["C2"] = "Actual (₹)"
        _apply_excel_styles(ws3, [ws3["A2"], ws3["B2"], ws3["C2"]])

        for r, cat in enumerate(categories, 3):
            ws3.cell(row=r, column=1).value = cat
            ws3.cell(row=r, column=2).value = 0
            ws3.cell(row=r, column=3).value = 0
            for c in [2, 3]:
                ws3.cell(row=r, column=c).number_format = '₹#,##0.00'

        for c in ["A", "B", "C"]:
            ws3.column_dimensions[c].width = 25

        # Save
        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}_budget.xlsx")
        wb.save(filepath)

        log_info(f"Budget Excel saved: {filepath}")
        _open_file(filepath)

        speak(f"Budget Excel sheet created Boss! Opening now.")
        return f"Budget Excel sheet created Boss! 3 sheets: Monthly Budget, Expense Log, Summary. Saved to Desktop."

    except ImportError:
        return "openpyxl not installed Boss. Run: pip install openpyxl"
    except Exception as e:
        log_error(f"Excel error: {e}")
        return f"Could not create Excel Boss: {e}"


def _create_tracker_excel(topic: str) -> str:
    """Create a habit/goal tracker Excel."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        from openpyxl.utils import get_column_letter
        import calendar

        wb  = openpyxl.Workbook()
        ws  = wb.active
        ws.title = "Tracker"

        now  = datetime.now()
        year = now.year
        month= now.month
        days = calendar.monthrange(year, month)[1]

        # Title
        ws.merge_cells(f"A1:{get_column_letter(days+3)}1")
        ws["A1"] = f"{topic.upper()} TRACKER — {now.strftime('%B %Y')}"
        ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws["A1"].fill = PatternFill("solid", fgColor="006600")
        ws["A1"].alignment = Alignment(horizontal="center")
        ws.row_dimensions[1].height = 25

        # Headers
        ws["A2"] = "Habit / Goal"
        ws["B2"] = "Target"
        for d in range(1, days+1):
            ws.cell(row=2, column=d+2).value = str(d)

        ws.cell(row=2, column=days+3).value = "Score"

        header_cells = [ws.cell(row=2, column=c) for c in range(1, days+4)]
        _apply_excel_styles(ws, header_cells, header_color="006600")

        # Sample habits
        habits = [
            ("Morning Exercise", "Daily"),
            ("Read for 30 mins", "Daily"),
            ("Drink 8 glasses water", "Daily"),
            ("Study Cybersecurity", "Daily"),
            ("Work on Cracka AI", "Daily"),
            ("Meditation", "Daily"),
            ("No junk food", "Daily"),
            (topic.title() if topic else "Custom Goal", "Daily"),
        ]

        for r, (habit, target) in enumerate(habits, 3):
            ws.cell(row=r, column=1).value = habit
            ws.cell(row=r, column=1).font = Font(name="Calibri", size=10)
            ws.cell(row=r, column=2).value = target

            # Score formula
            score_col = days + 3
            score_cell = ws.cell(row=r, column=score_col)
            data_cols = [get_column_letter(c) for c in range(3, days+3)]
            formula = "+".join([f"IF({col}{r}=\"✓\",1,0)" for col in data_cols])
            score_cell.value = f"={formula}"
            score_cell.font  = Font(bold=True)

        ws.column_dimensions["A"].width = 25
        ws.column_dimensions["B"].width = 10
        for d in range(1, days+1):
            ws.column_dimensions[get_column_letter(d+2)].width = 4
        ws.column_dimensions[get_column_letter(days+3)].width = 8

        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}_tracker.xlsx")
        wb.save(filepath)
        _open_file(filepath)

        speak(f"Tracker sheet created Boss!")
        return f"Habit tracker created Boss! Type ✓ for completed days. Saved to Desktop."

    except Exception as e:
        log_error(f"Tracker Excel error: {e}")
        return f"Could not create tracker Boss: {e}"


def _create_schedule_excel(topic: str) -> str:
    """Create a weekly schedule/timetable."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Weekly Schedule"

        days  = ["Monday", "Tuesday", "Wednesday", "Thursday",
                 "Friday", "Saturday", "Sunday"]
        times = [
            "6:00 AM", "7:00 AM", "8:00 AM", "9:00 AM", "10:00 AM",
            "11:00 AM", "12:00 PM", "1:00 PM", "2:00 PM", "3:00 PM",
            "4:00 PM", "5:00 PM", "6:00 PM", "7:00 PM", "8:00 PM",
            "9:00 PM", "10:00 PM"
        ]

        # Title
        ws.merge_cells("A1:H1")
        ws["A1"] = f"WEEKLY SCHEDULE — {topic.upper()}"
        ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws["A1"].fill = PatternFill("solid", fgColor="AA5500")
        ws["A1"].alignment = Alignment(horizontal="center")

        # Day headers
        ws["A2"] = "Time"
        for c, day in enumerate(days, 2):
            ws.cell(row=2, column=c).value = day
        _apply_excel_styles(
            ws,
            [ws.cell(row=2, column=c) for c in range(1, 9)],
            header_color="AA5500"
        )

        # Time rows
        for r, time in enumerate(times, 3):
            ws.cell(row=r, column=1).value = time
            ws.cell(row=r, column=1).font = Font(bold=True, size=10)
            ws.row_dimensions[r].height = 20

        ws.column_dimensions["A"].width = 12
        for c in range(2, 9):
            ws.column_dimensions[get_column_letter(c)].width = 18

        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}_schedule.xlsx")
        wb.save(filepath)
        _open_file(filepath)

        speak("Schedule created Boss!")
        return "Weekly schedule Excel created Boss! Fill in your activities. Saved to Desktop."

    except Exception as e:
        return f"Could not create schedule Boss: {e}"


def _create_student_excel(topic: str) -> str:
    """Create student marks/attendance tracker."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        from openpyxl.utils import get_column_letter

        wb  = openpyxl.Workbook()
        ws  = wb.active
        ws.title = "Student Records"

        subjects = [
            "Mathematics", "Physics", "Chemistry",
            "Computer Science", "English", "Physical Education"
        ]

        ws.merge_cells("A1:J1")
        ws["A1"] = f"STUDENT RECORDS — {topic.upper()}"
        ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws["A1"].fill = PatternFill("solid", fgColor="550077")
        ws["A1"].alignment = Alignment(horizontal="center")

        headers = ["Roll No", "Student Name"] + subjects + ["Total", "Percentage", "Grade"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c).value = h
        _apply_excel_styles(
            ws,
            [ws.cell(row=2, column=c) for c in range(1, len(headers)+1)],
            header_color="550077"
        )

        for r in range(3, 13):
            ws.cell(row=r, column=1).value = r - 2
            ws.cell(row=r, column=2).value = f"Student {r-2}"

            for c in range(3, len(subjects)+3):
                ws.cell(row=r, column=c).value = 0

            # Total formula
            sub_cols = [get_column_letter(c) for c in range(3, len(subjects)+3)]
            total_col = len(subjects) + 3
            pct_col   = total_col + 1
            grade_col = pct_col + 1

            total_cell = ws.cell(row=r, column=total_col)
            total_cell.value = f"=SUM({sub_cols[0]}{r}:{sub_cols[-1]}{r})"
            total_cell.font  = Font(bold=True)

            pct_cell = ws.cell(row=r, column=pct_col)
            pct_cell.value = f"={get_column_letter(total_col)}{r}/{len(subjects)}*100"
            pct_cell.number_format = "0.00"

            grade_cell = ws.cell(row=r, column=grade_col)
            pct_ref = f"{get_column_letter(pct_col)}{r}"
            grade_cell.value = (
                f'=IF({pct_ref}>=90,"A+",IF({pct_ref}>=80,"A",'
                f'IF({pct_ref}>=70,"B+",IF({pct_ref}>=60,"B",'
                f'IF({pct_ref}>=50,"C","F")))))'
            )

        ws.column_dimensions["A"].width = 8
        ws.column_dimensions["B"].width = 20
        for c in range(3, len(headers)+1):
            ws.column_dimensions[get_column_letter(c)].width = 14

        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}_students.xlsx")
        wb.save(filepath)
        _open_file(filepath)

        speak("Student records sheet created Boss!")
        return "Student records Excel created Boss! Grades calculate automatically. Saved to Desktop."

    except Exception as e:
        return f"Could not create student sheet Boss: {e}"


def _create_generic_excel(topic: str) -> str:
    """Create a generic data Excel sheet using AI."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Font, Alignment
        from openpyxl.utils import get_column_letter

        speak(f"Creating Excel sheet for {topic} Boss.")

        # Ask AI for column suggestions
        prompt = f"""For an Excel spreadsheet about "{topic}", suggest:
1. A list of 6-8 column headers
2. 3 sample data rows

Respond in this exact format:
COLUMNS: Col1, Col2, Col3, Col4, Col5, Col6
ROW1: val1, val2, val3, val4, val5, val6
ROW2: val1, val2, val3, val4, val5, val6
ROW3: val1, val2, val3, val4, val5, val6"""

        ai_response = ask_ai(prompt)

        columns  = []
        rows     = []

        for line in ai_response.split('\n'):
            line = line.strip()
            if line.startswith('COLUMNS:'):
                columns = [c.strip() for c in line[8:].split(',')]
            elif line.startswith('ROW'):
                vals = [v.strip() for v in line.split(':', 1)[-1].split(',')]
                rows.append(vals)

        if not columns:
            columns = ["ID", "Name", "Description", "Value", "Date", "Notes"]

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = topic[:30].title()

        # Title
        ws.merge_cells(f"A1:{get_column_letter(len(columns))}1")
        ws["A1"] = f"{topic.upper()} DATA SHEET"
        ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
        ws["A1"].fill = PatternFill("solid", fgColor="0077AA")
        ws["A1"].alignment = Alignment(horizontal="center")

        # Headers
        for c, col in enumerate(columns, 1):
            ws.cell(row=2, column=c).value = col
        _apply_excel_styles(ws, [ws.cell(row=2, column=c) for c in range(1, len(columns)+1)])

        # Sample data
        for r, row in enumerate(rows, 3):
            for c, val in enumerate(row[:len(columns)], 1):
                ws.cell(row=r, column=c).value = val

        for c in range(1, len(columns)+1):
            ws.column_dimensions[get_column_letter(c)].width = 18

        filename = _safe_filename(topic)
        filepath = os.path.join(DESKTOP, f"{filename}.xlsx")
        wb.save(filepath)
        _open_file(filepath)

        speak(f"Excel sheet for {topic} is ready Boss!")
        return f"Excel sheet created Boss! {len(columns)} columns ready. Saved to Desktop."

    except Exception as e:
        log_error(f"Generic Excel error: {e}")
        return f"Could not create Excel Boss: {e}"


# ══════════════════════════════════════════════════════════════
# COMMAND PARSER — called from ai_brain.py
# ══════════════════════════════════════════════════════════════

def handle_office_command(command: str) -> str:
    """
    Main entry point from ai_brain.py
    Detects which office app to use and calls appropriate function.
    """
    cmd = command.lower()

    # PowerPoint
    if any(w in cmd for w in [
        "presentation", "ppt", "slides", "powerpoint",
        "make slides", "create slides"
    ]):
        return create_presentation(command)

    # Word
    elif any(w in cmd for w in [
        "word document", "word doc", "report", "essay",
        "letter", "resume", "cv", "proposal", "write document",
        "create document"
    ]):
        return create_word_document(command)

    # Excel
    elif any(w in cmd for w in [
        "excel", "spreadsheet", "sheet", "tracker",
        "budget", "schedule", "timetable"
    ]):
        return create_excel_sheet(command)

    else:
        return ("I can create PowerPoint, Word documents, or Excel sheets Boss. "
                "Say 'make presentation on [topic]', "
                "'write report on [topic]', or "
                "'create excel sheet for [purpose]'.")